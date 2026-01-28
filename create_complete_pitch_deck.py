from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Color palette - Retro-Futuristic
COLOR_SURFACE = RGBColor(10, 17, 40)  # Dark navy background
COLOR_PRIMARY = RGBColor(0, 212, 255)  # Electric blue
COLOR_PRIMARY_FG = RGBColor(10, 17, 40)  # Dark navy for text on primary
COLOR_SECONDARY = RGBColor(255, 107, 53)  # Orange
COLOR_ACCENT = RGBColor(124, 58, 237)  # Purple
COLOR_MUTED = RGBColor(26, 40, 71)  # Lighter navy
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_MUTED_FG = RGBColor(148, 163, 184)  # Light gray

def add_background(slide):
    """Add dark navy background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_SURFACE

def add_slide_number(slide, number):
    """Add slide number in bottom right"""
    num_box = slide.shapes.add_textbox(Inches(9.2), Inches(5.2), Inches(0.6), Inches(0.3))
    num_frame = num_box.text_frame
    num_p = num_frame.paragraphs[0]
    num_p.alignment = PP_ALIGN.RIGHT
    num_run = num_p.add_run()
    num_run.text = str(number)
    num_run.font.size = Pt(12)
    num_run.font.color.rgb = COLOR_MUTED_FG

# SLIDE 1: Title Slide
def add_slide_1_title():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)

    # Left accent bar (orange)
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(3), Inches(5.625)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_SECONDARY
    left_bar.line.fill.background()

    # Logo/Brand
    title_box = slide.shapes.add_textbox(Inches(3.2), Inches(1.5), Inches(6.5), Inches(1))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = "StepTen.io"
    title_run.font.size = Pt(64)
    title_run.font.bold = True
    title_run.font.color.rgb = COLOR_PRIMARY

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(3.2), Inches(2.7), Inches(6.5), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_run = subtitle_p.add_run()
    subtitle_run.text = "The World's Most Advanced\nAI-Powered SEO Content Engine"
    subtitle_run.font.size = Pt(26)
    subtitle_run.font.color.rgb = COLOR_WHITE

    # Funding ask
    tag_box = slide.shapes.add_textbox(Inches(3.2), Inches(4.3), Inches(6.5), Inches(0.5))
    tag_frame = tag_box.text_frame
    tag_p = tag_frame.paragraphs[0]
    tag_run = tag_p.add_run()
    tag_run.text = "Raising $5M Series A"
    tag_run.font.size = Pt(22)
    tag_run.font.bold = True
    tag_run.font.color.rgb = COLOR_SECONDARY

# SLIDE 2: The Problem
def add_slide_2_problem():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 2)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "The Problem"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Problem boxes
    problems = [
        {"stat": "60%", "text": "of AI-generated content\nis penalized by Google"},
        {"stat": "8+ hrs", "text": "average time to create\none SEO-optimized article"},
        {"stat": "$89B", "text": "content marketing market\nwith no complete solution"}
    ]

    x_positions = [0.5, 3.5, 6.5]

    for i, prob in enumerate(problems):
        # Problem card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(1.8), Inches(2.7), Inches(3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_MUTED
        card.line.color.rgb = COLOR_PRIMARY
        card.line.width = Pt(2)

        # Stat
        stat_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.2), Inches(2.1), Inches(2.3), Inches(1)
        )
        stat_frame = stat_box.text_frame
        stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        stat_p = stat_frame.paragraphs[0]
        stat_p.alignment = PP_ALIGN.CENTER
        stat_run = stat_p.add_run()
        stat_run.text = prob["stat"]
        stat_run.font.size = Pt(56)
        stat_run.font.bold = True
        stat_run.font.color.rgb = COLOR_SECONDARY

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.2), Inches(3.3), Inches(2.3), Inches(1.2)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER
        desc_run = desc_p.add_run()
        desc_run.text = prob["text"]
        desc_run.font.size = Pt(16)
        desc_run.font.color.rgb = COLOR_WHITE

# SLIDE 3: Our Solution
def add_slide_3_solution():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 3)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Our Solution"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Subheader
    subheader_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    subheader_frame = subheader_box.text_frame
    subheader_p = subheader_frame.paragraphs[0]
    subheader_run = subheader_p.add_run()
    subheader_run.text = "8-Step AI Orchestration Pipeline with 5 Leading AI Models"
    subheader_run.font.size = Pt(20)
    subheader_run.font.color.rgb = COLOR_WHITE

    # 8 Steps
    steps = [
        "1. SERP Analysis", "2. Research & Data", "3. Outline Generation", "4. Content Creation",
        "5. SEO Optimization", "6. Quality Check", "7. Internal Linking", "8. Publish & Track"
    ]

    for i, step in enumerate(steps):
        row = i // 4
        col = i % 4

        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + col * 2.3), Inches(2 + row * 0.9), Inches(2.1), Inches(0.7)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = COLOR_ACCENT
        step_box.line.fill.background()

        text_box = slide.shapes.add_textbox(
            Inches(0.6 + col * 2.3), Inches(2.1 + row * 0.9), Inches(1.9), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.CENTER
        text_run = text_p.add_run()
        text_run.text = step
        text_run.font.size = Pt(14)
        text_run.font.bold = True
        text_run.font.color.rgb = COLOR_WHITE

    # AI Models
    models_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.4))
    models_frame = models_box.text_frame
    models_frame.word_wrap = True

    models_p = models_frame.paragraphs[0]
    models_p.alignment = PP_ALIGN.CENTER

    title_run = models_p.add_run()
    title_run.text = "Powered by 5 Leading AI Models\n"
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = COLOR_SECONDARY

    models_run = models_p.add_run()
    models_run.text = "Claude • GPT-4 • Grok • Gemini • Perplexity"
    models_run.font.size = Pt(22)
    models_run.font.color.rgb = COLOR_PRIMARY

# SLIDE 4: The Secret Weapon
def add_slide_4_secret_weapon():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 4)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "The Secret Weapon"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Subheader
    subheader_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    subheader_frame = subheader_box.text_frame
    subheader_p = subheader_frame.paragraphs[0]
    subheader_run = subheader_p.add_run()
    subheader_run.text = "11 Advanced SEO Features No Competitor Can Match"
    subheader_run.font.size = Pt(20)
    subheader_run.font.color.rgb = COLOR_WHITE

    # Features
    features = [
        "SERP Analysis", "Internal Linking", "Content Clusters",
        "NLP Entity Coverage", "Snippet Optimizer", "Freshness Detector",
        "27+ Languages", "A/B Testing", "Rank Tracking",
        "Score Predictor", "Autonomous AI Agent"
    ]

    for i, feature in enumerate(features):
        row = i // 4
        col = i % 4

        # Feature card
        if i < 8:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5 + col * 2.3), Inches(1.8 + row * 0.75), Inches(2.1), Inches(0.6)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_MUTED
            card.line.color.rgb = COLOR_PRIMARY
            card.line.width = Pt(1.5)

            text_box = slide.shapes.add_textbox(
                Inches(0.6 + col * 2.3), Inches(1.9 + row * 0.75), Inches(1.9), Inches(0.4)
            )
        else:
            # Last 3 features in bottom row
            col_offset = (i - 8) * 2.3 + 1.9
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(col_offset), Inches(4.3), Inches(2.1), Inches(0.6)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_SECONDARY
            card.line.fill.background()

            text_box = slide.shapes.add_textbox(
                Inches(col_offset + 0.1), Inches(4.4), Inches(1.9), Inches(0.4)
            )

        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.CENTER
        text_run = text_p.add_run()
        text_run.text = feature
        text_run.font.size = Pt(13)
        text_run.font.bold = True
        text_run.font.color.rgb = COLOR_WHITE

# SLIDE 5: Competitive Advantage
def add_slide_5_competitive():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 5)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Competitive Advantage"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Table header
    headers = ["Feature", "StepTen.io", "Jasper/Copy.ai", "SurferSEO", "Ahrefs"]
    col_widths = [2.5, 1.5, 1.5, 1.5, 1.5]
    x_start = 0.5
    y_start = 1.3

    for i, header in enumerate(headers):
        x_pos = sum(col_widths[:i]) + x_start

        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos), Inches(y_start), Inches(col_widths[i]), Inches(0.5)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_PRIMARY if i == 1 else COLOR_MUTED
        header_shape.line.fill.background()

        text_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.1), Inches(y_start + 0.05), Inches(col_widths[i] - 0.2), Inches(0.4)
        )
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.CENTER
        text_run = text_p.add_run()
        text_run.text = header
        text_run.font.size = Pt(12)
        text_run.font.bold = True
        text_run.font.color.rgb = COLOR_PRIMARY_FG if i == 1 else COLOR_WHITE

    # Table rows
    rows = [
        ["Multi-AI Orchestration", "✓", "✗", "✗", "✗"],
        ["SERP Analysis", "✓", "✗", "✓", "✓"],
        ["Content Clusters", "✓", "✗", "✗", "✗"],
        ["Internal Linking AI", "✓", "✗", "✗", "✗"],
        ["27+ Languages", "✓", "Limited", "Limited", "✗"],
        ["Autonomous Agent", "✓", "✗", "✗", "✗"],
        ["Rank Tracking", "✓", "✗", "✓", "✓"],
        ["A/B Testing", "✓", "✗", "✗", "✗"]
    ]

    row_height = 0.4
    for row_idx, row in enumerate(rows):
        y_pos = y_start + 0.5 + (row_idx * row_height)

        for col_idx, cell in enumerate(row):
            x_pos = sum(col_widths[:col_idx]) + x_start

            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(y_pos), Inches(col_widths[col_idx]), Inches(row_height)
            )
            cell_shape.fill.solid()
            cell_shape.fill.fore_color.rgb = COLOR_MUTED if row_idx % 2 == 0 else COLOR_SURFACE
            cell_shape.line.color.rgb = COLOR_PRIMARY if col_idx == 1 else COLOR_MUTED
            cell_shape.line.width = Pt(0.5)

            text_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(y_pos + 0.05),
                Inches(col_widths[col_idx] - 0.2), Inches(row_height - 0.1)
            )
            text_frame = text_box.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_p = text_frame.paragraphs[0]
            text_p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            text_run = text_p.add_run()
            text_run.text = cell
            text_run.font.size = Pt(11)
            text_run.font.color.rgb = COLOR_PRIMARY if (col_idx == 1 and cell == "✓") else COLOR_WHITE

# SLIDE 6: Market Opportunity
def add_slide_6_market():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 6)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Market Opportunity"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # TAM/SAM/SOM circles
    markets = [
        {"label": "TAM", "value": "$89B", "subtitle": "Content Marketing", "color": COLOR_ACCENT},
        {"label": "SAM", "value": "$24B", "subtitle": "SEO Software", "color": COLOR_PRIMARY},
        {"label": "Target", "value": "3 Segments", "subtitle": "Agencies • SaaS • Enterprise", "color": COLOR_SECONDARY}
    ]

    x_positions = [1, 3.75, 6.5]

    for i, market in enumerate(markets):
        # Circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_positions[i]), Inches(1.5), Inches(2.5), Inches(2.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = market["color"]
        circle.line.fill.background()

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.3), Inches(2), Inches(1.9), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        label_p = label_frame.paragraphs[0]
        label_p.alignment = PP_ALIGN.CENTER
        label_run = label_p.add_run()
        label_run.text = market["label"]
        label_run.font.size = Pt(18)
        label_run.font.bold = True
        label_run.font.color.rgb = COLOR_WHITE

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.3), Inches(2.5), Inches(1.9), Inches(0.6)
        )
        value_frame = value_box.text_frame
        value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        value_p = value_frame.paragraphs[0]
        value_p.alignment = PP_ALIGN.CENTER
        value_run = value_p.add_run()
        value_run.text = market["value"]
        value_run.font.size = Pt(32)
        value_run.font.bold = True
        value_run.font.color.rgb = COLOR_WHITE

        # Subtitle
        sub_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.1), Inches(3.2), Inches(2.3), Inches(0.6)
        )
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        sub_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        sub_p = sub_frame.paragraphs[0]
        sub_p.alignment = PP_ALIGN.CENTER
        sub_run = sub_p.add_run()
        sub_run.text = market["subtitle"]
        sub_run.font.size = Pt(12)
        sub_run.font.color.rgb = COLOR_WHITE

    # Target segments detail
    segments_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    segments_frame = segments_box.text_frame
    segments_p = segments_frame.paragraphs[0]
    segments_p.alignment = PP_ALIGN.CENTER

    seg_title = segments_p.add_run()
    seg_title.text = "Primary Target Segments:\n"
    seg_title.font.size = Pt(14)
    seg_title.font.bold = True
    seg_title.font.color.rgb = COLOR_MUTED_FG

    seg_detail = segments_p.add_run()
    seg_detail.text = "SEO Agencies (40%) • SaaS Companies (35%) • Enterprise Marketing Teams (25%)"
    seg_detail.font.size = Pt(16)
    seg_detail.font.color.rgb = COLOR_WHITE

# SLIDE 7: Business Model
def add_slide_7_business_model():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 7)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Business Model"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Pricing tiers
    tiers = [
        {"name": "Starter", "price": "$299", "features": "50 articles/mo\n5 projects\nBasic SEO", "color": COLOR_MUTED},
        {"name": "Professional", "price": "$999", "features": "200 articles/mo\nUnlimited projects\nAdvanced SEO", "color": COLOR_ACCENT},
        {"name": "Enterprise", "price": "Custom", "features": "Unlimited articles\nWhite label\nDedicated support", "color": COLOR_SECONDARY}
    ]

    x_positions = [0.5, 3.5, 6.5]

    for i, tier in enumerate(tiers):
        # Tier card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(1.3), Inches(2.7), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tier["color"]
        card.line.fill.background()

        # Tier name
        name_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.2), Inches(1.5), Inches(2.3), Inches(0.4)
        )
        name_frame = name_box.text_frame
        name_p = name_frame.paragraphs[0]
        name_p.alignment = PP_ALIGN.CENTER
        name_run = name_p.add_run()
        name_run.text = tier["name"]
        name_run.font.size = Pt(20)
        name_run.font.bold = True
        name_run.font.color.rgb = COLOR_WHITE

        # Price
        price_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.2), Inches(2), Inches(2.3), Inches(0.6)
        )
        price_frame = price_box.text_frame
        price_p = price_frame.paragraphs[0]
        price_p.alignment = PP_ALIGN.CENTER
        price_run = price_p.add_run()
        price_run.text = tier["price"]
        price_run.font.size = Pt(36)
        price_run.font.bold = True
        price_run.font.color.rgb = COLOR_WHITE

        # Per month
        if tier["price"] != "Custom":
            pm_run = price_p.add_run()
            pm_run.text = "/mo"
            pm_run.font.size = Pt(14)
            pm_run.font.color.rgb = COLOR_WHITE

        # Features
        feat_box = slide.shapes.add_textbox(
            Inches(x_positions[i] + 0.2), Inches(2.8), Inches(2.3), Inches(0.8)
        )
        feat_frame = feat_box.text_frame
        feat_frame.word_wrap = True
        feat_p = feat_frame.paragraphs[0]
        feat_p.alignment = PP_ALIGN.CENTER
        feat_run = feat_p.add_run()
        feat_run.text = tier["features"]
        feat_run.font.size = Pt(12)
        feat_run.font.color.rgb = COLOR_WHITE

    # Key metrics
    metrics = [
        {"label": "ACV", "value": "$25K"},
        {"label": "Gross Margin", "value": "85%"},
        {"label": "NDR", "value": "120%"}
    ]

    x_metric_positions = [1.5, 4, 6.5]

    for i, metric in enumerate(metrics):
        metric_box = slide.shapes.add_textbox(
            Inches(x_metric_positions[i]), Inches(4.2), Inches(2), Inches(1)
        )
        metric_frame = metric_box.text_frame
        metric_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        metric_p = metric_frame.paragraphs[0]
        metric_p.alignment = PP_ALIGN.CENTER

        label_run = metric_p.add_run()
        label_run.text = metric["label"] + "\n"
        label_run.font.size = Pt(14)
        label_run.font.color.rgb = COLOR_MUTED_FG

        value_run = metric_p.add_run()
        value_run.text = metric["value"]
        value_run.font.size = Pt(28)
        value_run.font.bold = True
        value_run.font.color.rgb = COLOR_PRIMARY

# SLIDE 8: Traction
def add_slide_8_traction():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 8)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Traction"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Status badge
    status_badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.7)
    )
    status_badge.fill.solid()
    status_badge.fill.fore_color.rgb = COLOR_SECONDARY
    status_badge.line.fill.background()

    status_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(0.5))
    status_frame = status_box.text_frame
    status_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    status_p = status_frame.paragraphs[0]
    status_p.alignment = PP_ALIGN.CENTER
    status_run = status_p.add_run()
    status_run.text = "100% Production Ready"
    status_run.font.size = Pt(28)
    status_run.font.bold = True
    status_run.font.color.rgb = COLOR_WHITE

    # Metrics
    metrics = [
        {"stat": "31", "label": "Database Tables"},
        {"stat": "18", "label": "API Endpoints"},
        {"stat": "15K+", "label": "Lines of Code"},
        {"stat": "5", "label": "AI Models Integrated"},
        {"stat": "27+", "label": "Languages Supported"},
        {"stat": "11", "label": "Advanced SEO Features"}
    ]

    for i, metric in enumerate(metrics):
        row = i // 3
        col = i % 3

        x_pos = 0.5 + col * 3.1
        y_pos = 2.5 + row * 1.4

        # Metric card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos), Inches(2.8), Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_MUTED
        card.line.color.rgb = COLOR_PRIMARY
        card.line.width = Pt(2)

        # Stat
        stat_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.2), Inches(y_pos + 0.1), Inches(2.4), Inches(0.6)
        )
        stat_frame = stat_box.text_frame
        stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        stat_p = stat_frame.paragraphs[0]
        stat_p.alignment = PP_ALIGN.CENTER
        stat_run = stat_p.add_run()
        stat_run.text = metric["stat"]
        stat_run.font.size = Pt(40)
        stat_run.font.bold = True
        stat_run.font.color.rgb = COLOR_PRIMARY

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.2), Inches(y_pos + 0.7), Inches(2.4), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        label_p = label_frame.paragraphs[0]
        label_p.alignment = PP_ALIGN.CENTER
        label_run = label_p.add_run()
        label_run.text = metric["label"]
        label_run.font.size = Pt(14)
        label_run.font.color.rgb = COLOR_WHITE

# SLIDE 9: Use of Funds
def add_slide_9_use_of_funds():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 9)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Use of Funds"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Subheader
    subheader_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    subheader_frame = subheader_box.text_frame
    subheader_p = subheader_frame.paragraphs[0]
    subheader_p.alignment = PP_ALIGN.CENTER
    subheader_run = subheader_p.add_run()
    subheader_run.text = "$5M Raising • 18 Month Runway"
    subheader_run.font.size = Pt(20)
    subheader_run.font.bold = True
    subheader_run.font.color.rgb = COLOR_SECONDARY

    # Allocation bars
    allocations = [
        {"category": "Product Development", "percent": 40, "amount": "$2.0M", "color": COLOR_PRIMARY},
        {"category": "Sales & Marketing", "percent": 35, "amount": "$1.75M", "color": COLOR_SECONDARY},
        {"category": "AI Infrastructure", "percent": 15, "amount": "$750K", "color": COLOR_ACCENT},
        {"category": "Operations & Legal", "percent": 10, "amount": "$500K", "color": COLOR_MUTED}
    ]

    y_start = 2.2
    bar_height = 0.7

    for i, alloc in enumerate(allocations):
        y_pos = y_start + i * 0.9

        # Background bar
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y_pos), Inches(6.5), Inches(bar_height)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = COLOR_MUTED
        bg_bar.line.fill.background()

        # Filled bar (proportional to percentage)
        filled_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y_pos), Inches(6.5 * alloc["percent"] / 100), Inches(bar_height)
        )
        filled_bar.fill.solid()
        filled_bar.fill.fore_color.rgb = alloc["color"]
        filled_bar.line.fill.background()

        # Category label
        cat_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos + 0.1), Inches(1.8), Inches(0.5)
        )
        cat_frame = cat_box.text_frame
        cat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cat_p = cat_frame.paragraphs[0]
        cat_p.alignment = PP_ALIGN.RIGHT
        cat_run = cat_p.add_run()
        cat_run.text = alloc["category"]
        cat_run.font.size = Pt(14)
        cat_run.font.bold = True
        cat_run.font.color.rgb = COLOR_WHITE

        # Percentage on bar
        perc_box = slide.shapes.add_textbox(
            Inches(2.6), Inches(y_pos + 0.15), Inches(1), Inches(0.4)
        )
        perc_frame = perc_box.text_frame
        perc_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        perc_p = perc_frame.paragraphs[0]
        perc_run = perc_p.add_run()
        perc_run.text = f"{alloc['percent']}%"
        perc_run.font.size = Pt(18)
        perc_run.font.bold = True
        perc_run.font.color.rgb = COLOR_WHITE

        # Amount
        amt_box = slide.shapes.add_textbox(
            Inches(7.8), Inches(y_pos + 0.15), Inches(1), Inches(0.4)
        )
        amt_frame = amt_box.text_frame
        amt_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        amt_p = amt_frame.paragraphs[0]
        amt_p.alignment = PP_ALIGN.RIGHT
        amt_run = amt_p.add_run()
        amt_run.text = alloc["amount"]
        amt_run.font.size = Pt(16)
        amt_run.font.bold = True
        amt_run.font.color.rgb = alloc["color"]

# SLIDE 10: Why Now
def add_slide_10_why_now():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 10)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Why Now?"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Reasons
    reasons = [
        {
            "title": "AI Convergence",
            "desc": "Multiple AI models now available with distinct strengths - orchestration is the new competitive advantage",
            "color": COLOR_PRIMARY
        },
        {
            "title": "Google's AI Crackdown",
            "desc": "March 2024 update penalizing low-quality AI content - massive need for better solutions",
            "color": COLOR_SECONDARY
        },
        {
            "title": "Content Arms Race",
            "desc": "Companies need 10x more content to compete - manual processes can't scale",
            "color": COLOR_ACCENT
        },
        {
            "title": "Market Fragmentation",
            "desc": "Existing tools are point solutions - no one offers complete end-to-end workflow",
            "color": COLOR_MUTED
        }
    ]

    for i, reason in enumerate(reasons):
        row = i // 2
        col = i % 2

        x_pos = 0.5 + col * 4.7
        y_pos = 1.5 + row * 1.8

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos), Inches(4.4), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = reason["color"]
        card.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.2), Inches(y_pos + 0.15), Inches(4), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_run = title_p.add_run()
        title_run.text = reason["title"]
        title_run.font.size = Pt(20)
        title_run.font.bold = True
        title_run.font.color.rgb = COLOR_WHITE

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.2), Inches(y_pos + 0.6), Inches(4), Inches(0.8)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_run = desc_p.add_run()
        desc_run.text = reason["desc"]
        desc_run.font.size = Pt(13)
        desc_run.font.color.rgb = COLOR_WHITE

# SLIDE 11: The Team
def add_slide_11_team():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)
    add_slide_number(slide, 11)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "The Team"
    header_run.font.size = Pt(44)
    header_run.font.bold = True
    header_run.font.color.rgb = COLOR_PRIMARY

    # Team member card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(1.5), Inches(6), Inches(3.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_MUTED
    card.line.color.rgb = COLOR_PRIMARY
    card.line.width = Pt(3)

    # Avatar placeholder
    avatar = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.25), Inches(1.8), Inches(1.5), Inches(1.5)
    )
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = COLOR_SECONDARY
    avatar.line.fill.background()

    # Initials
    init_box = slide.shapes.add_textbox(Inches(4.25), Inches(2.1), Inches(1.5), Inches(0.9))
    init_frame = init_box.text_frame
    init_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    init_p = init_frame.paragraphs[0]
    init_p.alignment = PP_ALIGN.CENTER
    init_run = init_p.add_run()
    init_run.text = "ST"
    init_run.font.size = Pt(48)
    init_run.font.bold = True
    init_run.font.color.rgb = COLOR_WHITE

    # Name
    name_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(5), Inches(0.5))
    name_frame = name_box.text_frame
    name_p = name_frame.paragraphs[0]
    name_p.alignment = PP_ALIGN.CENTER
    name_run = name_p.add_run()
    name_run.text = "Stephen Ten"
    name_run.font.size = Pt(32)
    name_run.font.bold = True
    name_run.font.color.rgb = COLOR_WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.1), Inches(5), Inches(0.4))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.add_run()
    title_run.text = "Founder & CEO"
    title_run.font.size = Pt(20)
    title_run.font.color.rgb = COLOR_SECONDARY

    # Bio
    bio_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.6), Inches(5), Inches(0.3))
    bio_frame = bio_box.text_frame
    bio_p = bio_frame.paragraphs[0]
    bio_p.alignment = PP_ALIGN.CENTER
    bio_run = bio_p.add_run()
    bio_run.text = "20+ years experience in SEO, AI, and SaaS development"
    bio_run.font.size = Pt(14)
    bio_run.font.color.rgb = COLOR_MUTED_FG

# SLIDE 12: Closing
def add_slide_12_closing():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide)

    # Right accent bar
    right_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7), Inches(0), Inches(3), Inches(5.625)
    )
    right_bar.fill.solid()
    right_bar.fill.fore_color.rgb = COLOR_PRIMARY
    right_bar.line.fill.background()

    # Main heading
    heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(1))
    heading_frame = heading_box.text_frame
    heading_p = heading_frame.paragraphs[0]
    heading_run = heading_p.add_run()
    heading_run.text = "Join Us"
    heading_run.font.size = Pt(56)
    heading_run.font.bold = True
    heading_run.font.color.rgb = COLOR_PRIMARY

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6), Inches(0.8))
    tag_frame = tag_box.text_frame
    tag_frame.word_wrap = True
    tag_p = tag_frame.paragraphs[0]
    tag_run = tag_p.add_run()
    tag_run.text = "Help us build the future of\nAI-powered content creation"
    tag_run.font.size = Pt(24)
    tag_run.font.color.rgb = COLOR_WHITE

    # Funding ask
    funding_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(6), Inches(0.6))
    funding_frame = funding_box.text_frame
    funding_p = funding_frame.paragraphs[0]
    funding_run = funding_p.add_run()
    funding_run.text = "Raising $5M Series A"
    funding_run.font.size = Pt(32)
    funding_run.font.bold = True
    funding_run.font.color.rgb = COLOR_SECONDARY

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(6), Inches(0.8))
    contact_frame = contact_box.text_frame
    contact_p = contact_frame.paragraphs[0]

    email_run = contact_p.add_run()
    email_run.text = "stephen@stepten.io\n"
    email_run.font.size = Pt(18)
    email_run.font.color.rgb = COLOR_WHITE

    web_run = contact_p.add_run()
    web_run.text = "www.stepten.io"
    web_run.font.size = Pt(18)
    web_run.font.color.rgb = COLOR_PRIMARY

# Generate all slides
print("Creating StepTen.io Investor Pitch Deck...")
add_slide_1_title()
print("✓ Slide 1: Title")
add_slide_2_problem()
print("✓ Slide 2: The Problem")
add_slide_3_solution()
print("✓ Slide 3: Our Solution")
add_slide_4_secret_weapon()
print("✓ Slide 4: Secret Weapon")
add_slide_5_competitive()
print("✓ Slide 5: Competitive Advantage")
add_slide_6_market()
print("✓ Slide 6: Market Opportunity")
add_slide_7_business_model()
print("✓ Slide 7: Business Model")
add_slide_8_traction()
print("✓ Slide 8: Traction")
add_slide_9_use_of_funds()
print("✓ Slide 9: Use of Funds")
add_slide_10_why_now()
print("✓ Slide 10: Why Now")
add_slide_11_team()
print("✓ Slide 11: The Team")
add_slide_12_closing()
print("✓ Slide 12: Closing")

# Save presentation
output_path = "/sessions/blissful-vibrant-lovelace/mnt/StepTen.io/StepTen_Investor_Pitch_Deck.pptx"
prs.save(output_path)
print(f"\n✓ Complete pitch deck created successfully!")
print(f"✓ Saved to: {output_path}")
print(f"✓ Total slides: 12")
